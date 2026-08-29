from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.api.coding_document_adapter_service import extract_document_preview
from app.api import coding_audit_service
from app.api.coding_document_export_service import apply_document_export, plan_document_export
from app.api.schemas.coding_documents import CodingDocumentExportApplyRequest as _CodingDocumentExportApplyRequest, CodingDocumentExportPlanRequest
from tests.coding_approval_test_helpers import approval_fields_for_plan


def CodingDocumentExportApplyRequest(**kwargs):
    plan = plan_document_export(
        CodingDocumentExportPlanRequest(
            **{key: value for key, value in kwargs.items() if key in {"session_id", "workspace_root", "file_path", "approval_granted", "approval_reason", "max_chars", "max_tables", "max_rows", "export_format", "target_path"}}
        )
    )
    approval = approval_fields_for_plan(workspace_root=kwargs["workspace_root"], operation_kind="document_export", mutation_class="document_export", source_file=kwargs["file_path"], plan=plan)
    return _CodingDocumentExportApplyRequest(**approval, **kwargs)


def test_document_export_plan_and_approved_execution(tmp_path: Path):
    path = tmp_path / "sample.docx"
    document = Document()
    document.add_paragraph("Export me safely")
    document.save(path)

    plan = plan_document_export(
        CodingDocumentExportPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            target_path="sample.export.md",
        )
    )

    assert plan.status == "planned"
    assert "Export me safely" in (plan.preview or "")

    result = apply_document_export(
        CodingDocumentExportApplyRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operator_approved=True,
            target_path="sample.export.md",
            expected_source_hash=plan.source_hash,
        )
    )

    assert result.status == "applied"
    assert result.mutation_performed is True
    assert (tmp_path / "sample.export.md").exists()


def _make_pdf(path: Path) -> None:
    import fitz

    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "PDF export text")
    document.save(path)
    document.close()


def _make_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["Name", "Value"])
    sheet.append(["Alpha", 1])
    workbook.save(path)


def _make_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "PPTX export text"
    presentation.save(path)


def _make_odt(path: Path) -> None:
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    document = OpenDocumentText()
    document.text.addElement(P(text="ODT export text"))
    document.save(str(path))


def _make_ods(path: Path) -> None:
    from odf.opendocument import OpenDocumentSpreadsheet
    from odf.table import Table, TableCell, TableRow
    from odf.text import P

    document = OpenDocumentSpreadsheet()
    sheet = Table(name="Data")
    row = TableRow()
    cell = TableCell()
    cell.addElement(P(text="ODS export text"))
    row.addElement(cell)
    sheet.addElement(row)
    document.spreadsheet.addElement(sheet)
    document.save(str(path))


def _make_odp(path: Path) -> None:
    from odf.draw import Frame, Page, TextBox
    from odf.opendocument import OpenDocumentPresentation
    from odf.text import P

    document = OpenDocumentPresentation()
    page = Page(masterpagename="Default", name="Slide 1")
    frame = Frame()
    box = TextBox()
    box.addElement(P(text="ODP export text"))
    frame.addElement(box)
    page.addElement(frame)
    document.presentation.addElement(page)
    document.save(str(path))


def test_export_quality_for_supported_document_formats(tmp_path: Path):
    docx = tmp_path / "sample.docx"
    document = Document()
    document.add_heading("DOCX export text", level=1)
    document.add_paragraph("Paragraph body")
    document.save(docx)
    makers = {
        "sample.pdf": _make_pdf,
        "sample.xlsx": _make_xlsx,
        "sample.pptx": _make_pptx,
        "sample.odt": _make_odt,
        "sample.ods": _make_ods,
        "sample.odp": _make_odp,
    }
    for filename, maker in makers.items():
        maker(tmp_path / filename)

    for path in [docx, *(tmp_path / name for name in makers)]:
        preview = extract_document_preview(path)
        plan = plan_document_export(
            CodingDocumentExportPlanRequest(
                workspace_root=str(tmp_path),
                file_path=str(path),
                approval_granted=True,
                export_format="markdown",
                target_path=f"{path.name}.export.md",
                max_chars=8000,
            )
        )
        output = plan.preview or ""
        assert plan.status == "planned"
        assert path.name in output
        assert preview.descriptor.label in output
        assert "Provenance" in output
        assert "Safety Notes" in output
        assert len(output.encode("utf-8")) <= 4000


def test_document_export_audit_does_not_store_full_extracted_text(tmp_path: Path, monkeypatch):
    audit_root = tmp_path / "audit"
    monkeypatch.setattr(coding_audit_service, "AUDIT_ROOT", audit_root)
    path = tmp_path / "sample.docx"
    secret_sentence = "This long extracted sentence must not be copied into audit records."
    document = Document()
    document.add_paragraph(secret_sentence)
    document.save(path)

    plan = plan_document_export(
        CodingDocumentExportPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            target_path="sample.export.md",
        )
    )
    result = apply_document_export(
        CodingDocumentExportApplyRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operator_approved=True,
            target_path="sample.export.md",
            expected_source_hash=plan.source_hash,
        )
    )

    assert result.status == "applied"
    audit_text = "\n".join(item.read_text(encoding="utf-8") for item in audit_root.glob("*.json"))
    assert secret_sentence not in audit_text
