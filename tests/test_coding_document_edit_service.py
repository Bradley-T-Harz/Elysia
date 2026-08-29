from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import load_workbook, Workbook
from pptx import Presentation

from app.api import coding_audit_service
from app.api.coding_document_edit_service import apply_document_edit, plan_document_edit
from app.api.schemas.coding_documents import CodingDocumentEditApplyRequest as _CodingDocumentEditApplyRequest, CodingDocumentEditPlanRequest
from tests.coding_approval_test_helpers import approval_fields_for_plan


def CodingDocumentEditApplyRequest(**kwargs):
    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            **{key: value for key, value in kwargs.items() if key in {"session_id", "workspace_root", "file_path", "approval_granted", "approval_reason", "max_chars", "max_tables", "max_rows", "operation", "parameters"}}
        )
    )
    approval = approval_fields_for_plan(workspace_root=kwargs["workspace_root"], operation_kind="document_edit", mutation_class="document_edit", source_file=kwargs["file_path"], plan=plan)
    return _CodingDocumentEditApplyRequest(**approval, **kwargs)


def _make_pdf(path: Path, text: str = "PDF text") -> None:
    import fitz

    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), text)
    document.save(path)
    document.close()


def _make_pdf_form(path: Path) -> None:
    import fitz

    document = fitz.open()
    page = document.new_page()
    text_widget = fitz.Widget()
    text_widget.field_name = "full_name"
    text_widget.field_type = fitz.PDF_WIDGET_TYPE_TEXT
    text_widget.rect = fitz.Rect(72, 72, 260, 94)
    text_widget.field_value = ""
    page.add_widget(text_widget)
    checkbox = fitz.Widget()
    checkbox.field_name = "agree"
    checkbox.field_type = fitz.PDF_WIDGET_TYPE_CHECKBOX
    checkbox.rect = fitz.Rect(72, 112, 90, 130)
    checkbox.field_value = False
    page.add_widget(checkbox)
    choice = fitz.Widget()
    choice.field_name = "color"
    choice.field_type = fitz.PDF_WIDGET_TYPE_COMBOBOX
    choice.rect = fitz.Rect(72, 150, 220, 174)
    choice.choice_values = ["red", "green"]
    choice.field_value = "red"
    page.add_widget(choice)
    document.save(path)
    document.close()


def _pdf_field_values(path: Path) -> dict[str, object]:
    import fitz

    document = fitz.open(path)
    values: dict[str, object] = {}
    try:
        for page in document:
            for widget in list(page.widgets() or []):
                values[str(widget.field_name)] = widget.field_value
    finally:
        document.close()
    return values


def test_pdf_arbitrary_inline_text_edit_is_refused_with_safe_alternatives(tmp_path: Path):
    path = tmp_path / "sample.pdf"
    _make_pdf(path)

    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operation="replace_text",
            parameters={"old_text": "a", "new_text": "b"},
        )
    )

    assert plan.status == "refused"
    assert plan.blocked_reason == "unstable_or_unsupported_document_edit"
    assert "Markdown/text export" in " ".join(plan.warnings)


def test_pdf_extract_pages_writes_derived_copy_with_approval(tmp_path: Path):
    path = tmp_path / "sample.pdf"
    _make_pdf(path, "First page")
    target = tmp_path / "sample.pages.pdf"

    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operation="extract_pages",
            parameters={"pages": [1], "target_path": str(target)},
        )
    )

    assert plan.status == "planned"
    assert plan.target_relative_path == "sample.pages.pdf"

    result = apply_document_edit(
        CodingDocumentEditApplyRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operator_approved=True,
            operation="extract_pages",
            parameters={"pages": [1], "target_path": str(target)},
            expected_source_hash=plan.source_hash,
        )
    )

    assert result.status == "applied"
    assert result.mutation_performed is True
    assert result.target_relative_path == "sample.pages.pdf"
    assert target.exists()
    assert result.new_hash != result.previous_hash


def test_pdf_metadata_update_writes_derived_copy_without_mutating_source(tmp_path: Path):
    path = tmp_path / "sample.pdf"
    _make_pdf(path, "Metadata source")
    before = path.read_bytes()
    target = tmp_path / "sample.metadata.pdf"

    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operation="update_metadata",
            parameters={"title": "Derived Title", "target_path": str(target)},
        )
    )
    result = apply_document_edit(
        CodingDocumentEditApplyRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operator_approved=True,
            operation="update_metadata",
            parameters={"title": "Derived Title", "target_path": str(target)},
            expected_source_hash=plan.source_hash,
        )
    )

    assert result.status == "applied"
    assert target.exists()
    assert path.read_bytes() == before


def test_pdf_form_fill_text_checkbox_and_select_writes_derived_copy(tmp_path: Path):
    path = tmp_path / "form.pdf"
    _make_pdf_form(path)
    before = path.read_bytes()
    target = tmp_path / "form.filled.pdf"

    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operation="fill_form_fields",
            parameters={
                "fields": {
                    "full_name": "the operator",
                    "agree": True,
                    "color": "green",
                },
                "target_path": str(target),
            },
        )
    )

    assert plan.status == "planned"
    assert plan.operation_details["blocked_count"] == 0

    result = apply_document_edit(
        CodingDocumentEditApplyRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operator_approved=True,
            operation="fill_form_fields",
            parameters={
                "fields": {
                    "full_name": "the operator",
                    "agree": True,
                    "color": "green",
                },
                "target_path": str(target),
            },
            expected_source_hash=plan.source_hash,
        )
    )

    assert result.status == "applied"
    assert result.operation_details["updated_count"] == 3
    assert target.exists()
    assert path.read_bytes() == before
    values = _pdf_field_values(target)
    assert values["full_name"] == "the operator"
    assert values["agree"] not in {False, "Off", None}
    assert values["color"] == "green"


def test_pdf_form_fill_rejects_unknown_field_name(tmp_path: Path):
    path = tmp_path / "form.pdf"
    _make_pdf_form(path)

    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operation="fill_form_fields",
            parameters={"fields": {"missing": "value"}, "target_path": str(tmp_path / "out.pdf")},
        )
    )

    assert plan.status == "blocked"
    assert plan.blocked_reason == "pdf_form_field_validation_failed"
    assert plan.operation_details["field_results"][0]["reason"] == "unknown_field_name"


def test_pdf_form_fill_rejects_invalid_select_option(tmp_path: Path):
    path = tmp_path / "form.pdf"
    _make_pdf_form(path)

    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operation="fill_form_fields",
            parameters={"fields": {"color": "blue"}, "target_path": str(tmp_path / "out.pdf")},
        )
    )

    assert plan.status == "blocked"
    assert plan.blocked_reason == "pdf_form_field_validation_failed"
    assert plan.operation_details["field_results"][0]["reason"] == "invalid_option"


def test_pdf_form_fill_audit_does_not_store_full_pdf_contents_or_field_values(tmp_path: Path, monkeypatch):
    audit_root = tmp_path / "audit"
    monkeypatch.setattr(coding_audit_service, "AUDIT_ROOT", audit_root)
    path = tmp_path / "form.pdf"
    _make_pdf_form(path)
    target = tmp_path / "form.filled.pdf"
    secret_value = "PRIVATE_FORM_VALUE_SHOULD_NOT_BE_IN_AUDIT"

    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operation="fill_form_fields",
            parameters={"fields": {"full_name": secret_value}, "target_path": str(target)},
        )
    )
    result = apply_document_edit(
        CodingDocumentEditApplyRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operator_approved=True,
            operation="fill_form_fields",
            parameters={"fields": {"full_name": secret_value}, "target_path": str(target)},
            expected_source_hash=plan.source_hash,
        )
    )

    assert result.status == "applied"
    audit_text = "\n".join(item.read_text(encoding="utf-8") for item in audit_root.glob("*.json"))
    assert secret_value not in audit_text
    assert "%PDF" not in audit_text


def test_docx_append_paragraph_requires_and_uses_approval(tmp_path: Path):
    path = tmp_path / "sample.docx"
    document = Document()
    document.add_paragraph("Before")
    document.save(path)

    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operation="append_paragraph",
            parameters={"text": "After"},
        )
    )

    assert plan.status == "planned"

    blocked = apply_document_edit(
        CodingDocumentEditApplyRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operator_approved=False,
            operation="append_paragraph",
            parameters={"text": "After"},
            expected_source_hash=plan.source_hash,
        )
    )
    assert blocked.status == "approval_required"

    applied = apply_document_edit(
        CodingDocumentEditApplyRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operator_approved=True,
            operation="append_paragraph",
            parameters={"text": "After"},
            expected_source_hash=plan.source_hash,
        )
    )
    assert applied.status == "applied"
    reopened = Document(str(path))
    assert "After" in [paragraph.text for paragraph in reopened.paragraphs]


def test_docx_replace_paragraph_round_trip(tmp_path: Path):
    path = tmp_path / "sample.docx"
    document = Document()
    document.add_paragraph("Before")
    document.save(path)

    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operation="replace_paragraph",
            parameters={"paragraph_index": 1, "text": "After"},
        )
    )
    result = apply_document_edit(
        CodingDocumentEditApplyRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operator_approved=True,
            operation="replace_paragraph",
            parameters={"paragraph_index": 1, "text": "After"},
            expected_source_hash=plan.source_hash,
        )
    )

    assert result.status == "applied"
    assert Document(str(path)).paragraphs[0].text == "After"


def test_xlsx_stable_operations_round_trip(tmp_path: Path):
    path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet["A1"] = "Name"
    workbook.save(path)

    operations = [
        ("set_cell", {"sheet": "Data", "cell": "B1", "value": "Value"}),
        ("append_row", {"sheet": "Data", "row": ["Alpha", 3]}),
        ("create_sheet", {"title": "Notes"}),
        ("rename_sheet", {"sheet": "Notes", "new_title": "Summary"}),
    ]
    for operation, parameters in operations:
        plan = plan_document_edit(
            CodingDocumentEditPlanRequest(
                workspace_root=str(tmp_path),
                file_path=str(path),
                approval_granted=True,
                operation=operation,
                parameters=parameters,
            )
        )
        result = apply_document_edit(
            CodingDocumentEditApplyRequest(
                workspace_root=str(tmp_path),
                file_path=str(path),
                approval_granted=True,
                operator_approved=True,
                operation=operation,
                parameters=parameters,
                expected_source_hash=plan.source_hash,
            )
        )
        assert result.status == "applied"

    reopened = load_workbook(str(path), data_only=False)
    assert reopened["Data"]["B1"].value == "Value"
    assert reopened["Data"]["A2"].value == "Alpha"
    assert "Summary" in reopened.sheetnames


def test_pptx_stable_operations_round_trip(tmp_path: Path):
    path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Before"
    presentation.save(path)

    for operation, parameters in [
        ("replace_text", {"old_text": "Before", "new_text": "After"}),
        ("append_slide", {"title": "Second", "body": "Body text"}),
    ]:
        plan = plan_document_edit(
            CodingDocumentEditPlanRequest(
                workspace_root=str(tmp_path),
                file_path=str(path),
                approval_granted=True,
                operation=operation,
                parameters=parameters,
            )
        )
        result = apply_document_edit(
            CodingDocumentEditApplyRequest(
                workspace_root=str(tmp_path),
                file_path=str(path),
                approval_granted=True,
                operator_approved=True,
                operation=operation,
                parameters=parameters,
                expected_source_hash=plan.source_hash,
            )
        )
        assert result.status == "applied"

    reopened = Presentation(str(path))
    text = "\n".join(shape.text for slide in reopened.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "After" in text
    assert "Second" in text


def test_odf_edits_are_refused(tmp_path: Path):
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    path = tmp_path / "sample.odt"
    document = OpenDocumentText()
    document.text.addElement(P(text="ODT text"))
    document.save(str(path))

    plan = plan_document_edit(
        CodingDocumentEditPlanRequest(
            workspace_root=str(tmp_path),
            file_path=str(path),
            approval_granted=True,
            operation="replace_paragraph",
            parameters={"paragraph_index": 1, "text": "After"},
        )
    )

    assert plan.status == "refused"
    assert plan.blocked_reason == "unstable_or_unsupported_document_edit"
