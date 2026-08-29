from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.api.coding_document_adapter_service import extract_document_preview


def test_docx_preview_extracts_paragraphs_and_tables(tmp_path: Path):
    path = tmp_path / "sample.docx"
    document = Document()
    document.add_heading("Quarterly Notes", level=1)
    document.add_paragraph("Revenue grew without exposing secrets.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    document.save(path)

    preview = extract_document_preview(path)

    assert preview.status == "completed"
    assert "Quarterly Notes" in preview.text_preview
    assert preview.tables
    assert preview.outline


def test_xlsx_preview_treats_formula_as_inert_text(tmp_path: Path):
    path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet["A1"] = "Value"
    sheet["A2"] = "=SUM(1,2)"
    workbook.save(path)

    preview = extract_document_preview(path)

    assert preview.status == "completed"
    assert preview.tables
    assert "Formula cells are reported as inert metadata" in " ".join(preview.warnings)
    assert "[formula redacted as inert text]" in preview.text_preview


def test_pptx_preview_extracts_slide_text(tmp_path: Path):
    path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Launch Plan"
    presentation.save(path)

    preview = extract_document_preview(path)

    assert preview.status == "completed"
    assert "Launch Plan" in preview.text_preview
    assert preview.metadata["slide_count"] == 1
