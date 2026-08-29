"""Stable document edit planning and approved execution."""

from __future__ import annotations

from hashlib import sha256
from pathlib import Path
from uuid import uuid4

from app.api.coding_audit_service import write_coding_audit_record
from app.api.coding_backup_service import create_coding_backup
from app.api.coding_document_adapter_service import inspect_document
from app.api.coding_document_type_registry import detect_document_type
from app.api.coding_pdf_adapter import discover_pdf_form_fields
from app.api.coding_path_guard_service import guard_workspace_path
from app.api.coding_operation_hash_service import operation_plan_hash
from app.api.coding_operation_service import consume_operation_approval
from app.api.schemas.coding_documents import (
    CodingDocumentApplyResponse,
    CodingDocumentEditApplyRequest,
    CodingDocumentEditPlanRequest,
    CodingDocumentPlanResponse,
)


def _hash_bytes(path: Path) -> str:
    digest = sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _normalize_checkbox_value(value: object, on_state: str | None) -> str:
    if isinstance(value, bool):
        return (on_state or "Yes") if value else "Off"
    text = str(value).strip()
    if text.lower() in {"true", "yes", "on", "checked", "1"}:
        return on_state or "Yes"
    if text.lower() in {"false", "no", "off", "unchecked", "0"}:
        return "Off"
    if text in {"Off", on_state}:
        return text
    raise ValueError("invalid_checkbox_value")


def _refusal(file_label: str, relative_path: str | None, reason: str, warnings: list[str] | None = None) -> CodingDocumentPlanResponse:
    safe_alternatives = []
    if file_label.lower().endswith(".pdf"):
        safe_alternatives = [
            "Use Markdown/text export for editable extracted content.",
            "Create a corrected DOCX/Markdown source copy when format fidelity matters.",
            "Use a derived PDF annotation/stamp/highlight operation for visible review notes.",
            "Use redaction rectangles plus replacement overlay when exact coordinates are approved.",
            "Rebuild a derived PDF from approved corrected content rather than arbitrary inline text surgery.",
        ]
    return CodingDocumentPlanResponse(
        status="refused",
        action="document_edit",
        file_label=file_label,
        relative_path=relative_path,
        blocked_reason=reason,
        plan_summary=(
            "This requested document edit is not a stable governed operation for this format. "
            + ("Nearest safe alternatives: " + " ".join(safe_alternatives) if safe_alternatives else "")
        ).strip(),
        warnings=(warnings or []) + safe_alternatives,
    )


def _validate_pdf_form_fill_plan(path: Path, parameters: dict[str, object]) -> dict[str, object]:
    requested = _requested_form_values(parameters)
    discovered = discover_pdf_form_fields(path, max_fields=500)
    by_name = {str(field.get("name")): field for field in discovered if field.get("name")}
    field_results: list[dict[str, object]] = []
    for field_name, value in requested.items():
        field = by_name.get(str(field_name))
        if not field:
            field_results.append(
                {
                    "field_name": str(field_name),
                    "status": "blocked",
                    "reason": "unknown_field_name",
                }
            )
            continue
        field_type = str(field.get("field_type") or "").lower()
        if field.get("read_only") is True:
            field_results.append(
                {
                    "field_name": str(field_name),
                    "status": "blocked",
                    "field_type": field_type,
                    "reason": "field_is_read_only",
                }
            )
            continue
        if field_type in {"combobox", "listbox"}:
            options = [str(option) for option in field.get("options") or []]
            if options and str(value) not in options:
                field_results.append(
                    {
                        "field_name": str(field_name),
                        "status": "blocked",
                        "field_type": field_type,
                        "reason": "invalid_option",
                        "options": options,
                    }
                )
                continue
        if field_type in {"checkbox", "radio"}:
            try:
                _normalize_checkbox_value(value, str(field.get("on_state") or "") or None)
            except Exception as exc:
                field_results.append(
                    {
                        "field_name": str(field_name),
                        "status": "blocked",
                        "field_type": field_type,
                        "reason": str(exc),
                    }
                )
                continue
        if field_type in {"signature", "unknown"}:
            field_results.append(
                {
                    "field_name": str(field_name),
                    "status": "blocked",
                    "field_type": field_type,
                    "reason": "unsupported_field_type",
                }
            )
            continue
        field_results.append(
            {
                "field_name": str(field_name),
                "status": "planned",
                "field_type": field_type,
                "value_preview": "[provided]",
            }
        )
    return {
        "field_results": field_results,
        "requested_count": len(requested),
        "discovered_count": len(discovered),
        "blocked_count": sum(1 for item in field_results if item["status"] == "blocked"),
    }


def plan_document_edit(payload: CodingDocumentEditPlanRequest) -> CodingDocumentPlanResponse:
    guarded = guard_workspace_path(workspace_root=payload.workspace_root, target_path=payload.file_path)
    file_label = guarded.target_path.name
    if not guarded.allowed:
        return CodingDocumentPlanResponse(
            status="blocked",
            action="document_edit",
            file_label=file_label,
            relative_path=guarded.relative_path,
            blocked_reason=guarded.reason,
            plan_summary="Document edit is blocked by workspace/path policy.",
        )
    if not payload.approval_granted:
        return CodingDocumentPlanResponse(
            status="approval_required",
            action="document_edit",
            file_label=file_label,
            relative_path=guarded.relative_path,
            blocked_reason="explicit_approval_required",
            plan_summary="Document edit planning requires approval before source content is parsed.",
        )
    descriptor = detect_document_type(guarded.target_path)
    preview = inspect_document(guarded.target_path)
    if preview.blocked_reason:
        return CodingDocumentPlanResponse(
            status="blocked",
            action="document_edit",
            file_label=file_label,
            relative_path=guarded.relative_path,
            blocked_reason=preview.blocked_reason,
            plan_summary="Document edit is blocked by document safety policy.",
            warnings=preview.warnings,
        )
    if not descriptor.editable or payload.operation not in descriptor.stable_edit_operations:
        return _refusal(
            file_label,
            guarded.relative_path,
            "unstable_or_unsupported_document_edit",
            list(descriptor.notes),
        )
    pdf_target_relative: str | None = None
    if descriptor.adapter == "pdf":
        target_guard, target_error = _guard_pdf_target(payload, guarded.target_path)
        if target_error or not target_guard or not target_guard.allowed:
            return CodingDocumentPlanResponse(
                status="blocked",
                action="document_edit",
                file_label=file_label,
                relative_path=guarded.relative_path,
                blocked_reason=target_error or (target_guard.reason if target_guard else "pdf_target_blocked"),
                plan_summary="PDF derived-copy operation is blocked by target path policy.",
                warnings=list(descriptor.notes),
            )
        pdf_target_relative = target_guard.relative_path
        if payload.operation == "fill_form_fields":
            details = _validate_pdf_form_fill_plan(guarded.target_path, payload.parameters)
            if not details["field_results"]:
                return CodingDocumentPlanResponse(
                    status="blocked",
                    action="document_edit",
                    file_label=file_label,
                    relative_path=guarded.relative_path,
                    blocked_reason="no_form_fields_requested",
                    plan_summary="PDF form fill requires at least one requested field value.",
                    operation_details=details,
                    warnings=list(descriptor.notes),
                )
            if details["blocked_count"]:
                return CodingDocumentPlanResponse(
                    status="blocked",
                    action="document_edit",
                    file_label=file_label,
                    relative_path=guarded.relative_path,
                    target_relative_path=pdf_target_relative,
                    blocked_reason="pdf_form_field_validation_failed",
                    plan_summary="PDF form fill is blocked until requested field names and values match discovered writable fields.",
                    operation_details=details,
                    warnings=list(descriptor.notes),
                )
    source_hash = preview.safety.byte_hash
    operation_details = (
        _validate_pdf_form_fill_plan(guarded.target_path, payload.parameters)
        if descriptor.adapter == "pdf" and payload.operation == "fill_form_fields"
        else {}
    )
    target_relative = pdf_target_relative if descriptor.adapter == "pdf" else guarded.relative_path
    plan_hash = operation_plan_hash(
        action="document_edit",
        source_relative_path=guarded.relative_path,
        target_relative_path=target_relative,
        source_hash=source_hash,
        details={"operation": payload.operation, "parameters": payload.parameters},
    )
    return CodingDocumentPlanResponse(
        status="planned",
        action="document_edit",
        file_label=file_label,
        relative_path=guarded.relative_path,
        target_relative_path=target_relative,
        plan_summary=(
            f"Plan approved derived-copy PDF operation '{payload.operation}'. Source PDF will not be modified."
            if descriptor.adapter == "pdf"
            else f"Plan stable {descriptor.label} operation '{payload.operation}'."
        ),
        source_hash=source_hash,
        plan_hash=plan_hash,
        preview=str(payload.parameters)[:2000],
        operation_details=operation_details,
        warnings=list(descriptor.notes)
        + [
            "Execution requires explicit approval and source hash match.",
            *(
                [
                    "PDF operations write a derived PDF copy unless explicitly blocked; source PDF remains unchanged.",
                    "Arbitrary inline sentence editing in a PDF is not attempted.",
                ]
                if descriptor.adapter == "pdf"
                else []
            ),
        ],
    )


def _apply_docx(path: Path, operation: str, parameters: dict[str, object]) -> None:
    from docx import Document

    document = Document(str(path))
    if operation == "append_paragraph":
        document.add_paragraph(str(parameters.get("text") or ""))
    elif operation == "replace_paragraph":
        index = int(parameters.get("paragraph_index") or 1) - 1
        text = str(parameters.get("text") or "")
        if index < 0 or index >= len(document.paragraphs):
            raise ValueError("paragraph_index_out_of_range")
        document.paragraphs[index].text = text
    else:
        raise ValueError("unsupported_docx_operation")
    document.save(str(path))


def _apply_xlsx(path: Path, operation: str, parameters: dict[str, object]) -> None:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=False, data_only=False)
    sheet_name = str(parameters.get("sheet") or workbook.worksheets[0].title)
    if operation == "create_sheet":
        title = str(parameters.get("title") or "New Sheet")
        if title in workbook.sheetnames:
            raise ValueError("sheet_already_exists")
        workbook.create_sheet(title=title)
    elif operation == "rename_sheet":
        new_title = str(parameters.get("new_title") or "")
        if not new_title:
            raise ValueError("missing_new_title")
        workbook[sheet_name].title = new_title
    elif operation == "set_cell":
        cell = str(parameters.get("cell") or "")
        if not cell:
            raise ValueError("missing_cell")
        value = parameters.get("value")
        if isinstance(value, str) and value.startswith("="):
            raise ValueError("formula_write_refused")
        workbook[sheet_name][cell] = value
    elif operation == "append_row":
        row = parameters.get("row")
        if not isinstance(row, list):
            raise ValueError("row_must_be_list")
        if any(isinstance(item, str) and item.startswith("=") for item in row):
            raise ValueError("formula_write_refused")
        workbook[sheet_name].append(row)
    else:
        raise ValueError("unsupported_xlsx_operation")
    workbook.save(str(path))


def _apply_pptx(path: Path, operation: str, parameters: dict[str, object]) -> None:
    from pptx import Presentation

    presentation = Presentation(str(path))
    if operation == "append_slide":
        title = str(parameters.get("title") or "New Slide")
        body = str(parameters.get("body") or "")
        layout = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape != title_shape:
                shape.text = body
                break
    elif operation == "replace_text":
        old_text = str(parameters.get("old_text") or "")
        new_text = str(parameters.get("new_text") or "")
        if not old_text:
            raise ValueError("missing_old_text")
        replaced = False
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and old_text in shape.text:
                    shape.text = shape.text.replace(old_text, new_text)
                    replaced = True
        if not replaced:
            raise ValueError("text_not_found")
    else:
        raise ValueError("unsupported_pptx_operation")
    presentation.save(str(path))


def _int_list(value: object, *, default: list[int]) -> list[int]:
    if isinstance(value, list):
        result = []
        for item in value:
            try:
                result.append(int(item))
            except (TypeError, ValueError):
                continue
        return result or default
    if isinstance(value, str):
        result = []
        for item in value.split(","):
            try:
                result.append(int(item.strip()))
            except ValueError:
                continue
        return result or default
    return default


def _pdf_target_relative_path(payload: CodingDocumentEditPlanRequest, source_path: Path) -> str | None:
    target = payload.parameters.get("target_path")
    if isinstance(target, str) and target.strip():
        return target.strip()
    return f"{source_path.stem}.{payload.operation}.derived.pdf"


def _guard_pdf_target(payload: CodingDocumentEditPlanRequest, source_path: Path):
    target_path = _pdf_target_relative_path(payload, source_path)
    guarded = guard_workspace_path(
        workspace_root=payload.workspace_root,
        target_path=target_path or f"{source_path.stem}.derived.pdf",
        require_existing=False,
        allow_directory=False,
    )
    if guarded.allowed and guarded.target_path.resolve() == source_path.resolve():
        return None, "pdf_target_must_be_derived_copy"
    if guarded.allowed and guarded.target_path.suffix.lower() != ".pdf":
        return None, "pdf_target_must_be_pdf"
    return guarded, None


def _apply_pdf_derived(
    source_path: Path,
    payload: CodingDocumentEditApplyRequest,
) -> tuple[Path, str | None, dict[str, object]]:
    from app.api.coding_pdf_worker_client import run_pdf_worker

    target_guard, target_error = _guard_pdf_target(payload, source_path)
    if target_error:
        raise ValueError(target_error)
    if not target_guard or not target_guard.allowed:
        raise ValueError(target_guard.reason if target_guard else "pdf_target_blocked")

    target_guard.target_path.parent.mkdir(parents=True, exist_ok=True)
    operation = payload.operation
    parameters = payload.parameters
    worker_parameters = dict(parameters)
    if operation == "merge_pdf":
        other_path = parameters.get("other_pdf_path")
        if not isinstance(other_path, str) or not other_path.strip():
            raise ValueError("other_pdf_path_required")
        other_guard = guard_workspace_path(workspace_root=payload.workspace_root, target_path=other_path,
                                           require_existing=True, allow_directory=False)
        if not other_guard.allowed:
            raise ValueError(other_guard.reason or "other_pdf_blocked")
        if other_guard.target_path.suffix.lower() != ".pdf":
            raise ValueError("other_file_must_be_pdf")
        worker_parameters["other_pdf_path"] = str(other_guard.target_path)
    result = run_pdf_worker({"action": "apply_derived_operation", "source_path": str(source_path),
                             "target_path": str(target_guard.target_path), "operation": operation,
                             "parameters": worker_parameters})
    operation_details = result.get("operation_details")
    if not isinstance(operation_details, dict):
        operation_details = {}

    return target_guard.target_path, target_guard.relative_path, operation_details


def _requested_form_values(parameters: dict[str, object]) -> dict[str, object]:
    values = parameters.get("fields")
    if isinstance(values, dict):
        return values
    values = parameters.get("field_values")
    if isinstance(values, dict):
        return values
    name = parameters.get("field_name")
    if isinstance(name, str) and name.strip():
        return {name.strip(): parameters.get("value")}
    return {}


def apply_document_edit(payload: CodingDocumentEditApplyRequest) -> CodingDocumentApplyResponse:
    plan = plan_document_edit(payload)
    if plan.status != "planned":
        return CodingDocumentApplyResponse(
            status=plan.status,
            action="document_edit",
            file_label=plan.file_label,
            relative_path=plan.relative_path,
            blocked_reason=plan.blocked_reason,
            warnings=plan.warnings,
        )
    if not payload.operator_approved:
        return CodingDocumentApplyResponse(
            status="approval_required",
            action="document_edit",
            file_label=plan.file_label,
            relative_path=plan.relative_path,
            blocked_reason="operator_approval_required",
            warnings=["Document edit execution requires explicit operator approval."],
        )
    if not payload.expected_source_hash:
        return CodingDocumentApplyResponse(
            status="blocked",
            action="document_edit",
            file_label=plan.file_label,
            relative_path=plan.relative_path,
            blocked_reason="expected_source_hash_required",
            warnings=["Document edits require the exact planned source hash."],
        )
    if payload.expected_source_hash != plan.source_hash:
        return CodingDocumentApplyResponse(
            status="blocked",
            action="document_edit",
            file_label=plan.file_label,
            relative_path=plan.relative_path,
            blocked_reason="source_hash_mismatch",
            warnings=["Re-inspect the document before applying edits."],
        )
    guarded = guard_workspace_path(workspace_root=payload.workspace_root, target_path=payload.file_path)
    if not guarded.allowed:
        return CodingDocumentApplyResponse(
            status="blocked",
            action="document_edit",
            file_label=plan.file_label,
            relative_path=guarded.relative_path,
            blocked_reason=guarded.reason,
        )
    descriptor = detect_document_type(guarded.target_path)
    previous_hash = _hash_bytes(guarded.target_path)
    pdf_target_relative: str | None = None
    pdf_target_path: Path | None = None
    pdf_operation_details: dict[str, object] = {}
    if descriptor.adapter == "pdf" and plan.target_relative_path:
        pdf_target_guard = guard_workspace_path(
            workspace_root=payload.workspace_root,
            target_path=plan.target_relative_path,
            require_existing=False,
            allow_directory=False,
        )
        if pdf_target_guard.target_path.exists():
            return CodingDocumentApplyResponse(
                status="blocked",
                action="document_edit",
                file_label=plan.file_label,
                relative_path=guarded.relative_path,
                target_relative_path=plan.target_relative_path,
                blocked_reason="target_exists",
                warnings=["Derived PDF operations never overwrite an existing target."],
            )
    exact_files = [payload.file_path]
    if plan.target_relative_path and plan.target_relative_path != plan.relative_path:
        exact_files.append(plan.target_relative_path)
    approval = consume_operation_approval(
        approval_id=payload.approval_id,
        approval_token=payload.approval_token,
        operation_kind="document_edit",
        workspace_root=payload.workspace_root,
        exact_files=exact_files,
        source_hash=plan.source_hash,
        plan_hash=plan.plan_hash or "",
        allowed_mutation_class="document_edit",
    )
    if not approval.allowed:
        return CodingDocumentApplyResponse(
            status="approval_required",
            action="document_edit",
            file_label=plan.file_label,
            relative_path=guarded.relative_path,
            target_relative_path=plan.target_relative_path,
            approval_id=payload.approval_id,
            blocked_reason=approval.reason,
            warnings=["A matching one-time document approval is required."],
        )
    backup = None
    if descriptor.adapter != "pdf":
        backup = create_coding_backup(
            workspace_root=guarded.workspace_root,
            source_path=guarded.target_path,
            source_relative_path=guarded.relative_path or guarded.target_path.name,
            operation_kind="document_edit",
            session_id=payload.session_id,
        )
    try:
        if descriptor.adapter == "docx":
            _apply_docx(guarded.target_path, payload.operation, payload.parameters)
        elif descriptor.adapter == "xlsx":
            _apply_xlsx(guarded.target_path, payload.operation, payload.parameters)
        elif descriptor.adapter == "pptx":
            _apply_pptx(guarded.target_path, payload.operation, payload.parameters)
        elif descriptor.adapter == "pdf":
            pdf_target_path, pdf_target_relative, pdf_operation_details = _apply_pdf_derived(guarded.target_path, payload)
        else:
            raise ValueError("unsupported_document_edit_adapter")
    except Exception as exc:
        return CodingDocumentApplyResponse(
            status="blocked",
            action="document_edit",
            file_label=plan.file_label,
            relative_path=guarded.relative_path,
            blocked_reason=str(exc),
            previous_hash=previous_hash,
            approval_id=payload.approval_id,
            backup_relative_path=backup.backup_relative_path if backup else None,
            rollback_receipt_id=backup.receipt_id if backup else None,
            warnings=["No document mutation was completed."],
        )
    new_hash = _hash_bytes(pdf_target_path if pdf_target_path else guarded.target_path)
    audit_written = write_coding_audit_record(
        "document_edit",
        uuid4().hex[:16],
        {
            "session_id": payload.session_id,
            "relative_path": guarded.relative_path,
            "document_type": descriptor.type_id,
            "operation": payload.operation,
            "previous_hash": previous_hash,
            "new_hash": new_hash,
            "operator_approved": True,
            "approval_id": payload.approval_id,
            "plan_hash": plan.plan_hash,
            "backup_relative_path": backup.backup_relative_path if backup else None,
        },
    )
    return CodingDocumentApplyResponse(
        status="applied",
        action="document_edit",
        file_label=plan.file_label,
        relative_path=guarded.relative_path,
        target_relative_path=pdf_target_relative if descriptor.adapter == "pdf" else plan.target_relative_path,
        mutation_performed=True,
        audit_written=audit_written,
        previous_hash=previous_hash,
        new_hash=new_hash,
        approval_id=payload.approval_id,
        backup_relative_path=backup.backup_relative_path if backup else None,
        rollback_receipt_id=backup.receipt_id if backup else None,
        operation_details=pdf_operation_details if descriptor.adapter == "pdf" else {},
        warnings=[
            (
                "Derived PDF copy was written locally; source PDF was not modified."
                if descriptor.adapter == "pdf"
                else "Stable document edit was applied locally without shell, macros, formula execution, or external services."
            )
        ],
        rollback_note=(
            "Delete or restore the derived PDF copy if this operation was not desired. Source PDF remains unchanged."
            if descriptor.adapter == "pdf"
            else f"Restore from {backup.backup_relative_path} using receipt {backup.receipt_id}."
        ),
    )


__all__ = ("apply_document_edit", "plan_document_edit")
