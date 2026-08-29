"""PPTX adapter for bounded slide text/table extraction."""

from __future__ import annotations

from pathlib import Path
from typing import Any


def extract_pptx_preview(path: Path, *, max_chars: int, max_tables: int, max_rows: int) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    props = presentation.core_properties
    metadata = {
        "slide_count": len(presentation.slides),
        "title": props.title,
        "author": props.author,
        "subject": props.subject,
        "created": props.created.isoformat() if props.created else None,
        "modified": props.modified.isoformat() if props.modified else None,
    }
    text_parts: list[str] = []
    outline: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    provenance: list[dict[str, Any]] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_texts: list[str] = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "text") and shape.text:
                text = str(shape.text)
                slide_texts.append(text)
                if not outline or len(outline) < 100:
                    outline.append({"kind": "slide_text", "slide": slide_index, "shape": shape_index, "text": text[:300]})
            if getattr(shape, "has_table", False) and len(tables) < max_tables:
                rows: list[list[str]] = []
                for row in shape.table.rows[:max_rows]:
                    rows.append([cell.text for cell in row.cells])
                tables.append({"kind": "table", "slide": slide_index, "shape": shape_index, "rows": rows, "row_count_previewed": len(rows)})
        if slide_texts and len("\n".join(text_parts)) < max_chars:
            joined = "\n".join(slide_texts)
            text_parts.append(f"[slide {slide_index}]\n{joined}")
            provenance.append({"kind": "slide", "slide": slide_index, "chars": len(joined)})

    return {
        "status": "completed",
        "metadata": metadata,
        "text_preview": "\n\n".join(text_parts)[:max_chars],
        "tables": tables,
        "outline": outline,
        "provenance": provenance,
        "warnings": [],
    }


__all__ = ("extract_pptx_preview",)
